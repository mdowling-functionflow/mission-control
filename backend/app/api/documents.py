"""API routes for the document workspace."""

from __future__ import annotations

from uuid import UUID

from fastapi import APIRouter, Depends, HTTPException, Query, status
from sqlmodel import SQLModel, col, select
from sqlmodel.ext.asyncio.session import AsyncSession

import httpx
import mimetypes
from pathlib import Path

from app.api.deps import ORG_MEMBER_DEP, SESSION_DEP, require_org_admin
from app.core.config import settings
from app.core.time import utcnow
from app.models.documents import Document
from app.models.executive_agents import ExecutiveAgent
from app.schemas.documents import DocumentCreate, DocumentRead, DocumentUpdate
from app.services.organizations import OrganizationContext

router = APIRouter(prefix="/documents", tags=["documents"])


async def _resolve_agent(
    session: AsyncSession, doc: Document, org_id: UUID,
) -> DocumentRead:
    """Build a DocumentRead with resolved agent display info."""
    read = DocumentRead.model_validate(doc)
    if doc.source_agent_id:
        agent = await ExecutiveAgent.objects.filter_by(
            id=doc.source_agent_id,
            organization_id=org_id,
        ).first(session)
        if agent:
            read.agent_display_name = agent.display_name
            read.agent_avatar_emoji = agent.avatar_emoji
    return read


@router.get("", response_model=list[DocumentRead])
async def list_documents(
    source_agent_id: UUID | None = Query(default=None),
    doc_type: str | None = Query(default=None),
    origin: str | None = Query(default=None),
    status_filter: str | None = Query(default=None, alias="status"),
    limit: int = Query(default=50, le=200),
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = ORG_MEMBER_DEP,
) -> list[DocumentRead]:
    stmt = (
        select(Document)
        .where(col(Document.organization_id) == ctx.organization.id)
        .where(col(Document.status) != "archived")
        .order_by(col(Document.updated_at).desc())
        .limit(limit)
    )
    if source_agent_id:
        stmt = stmt.where(col(Document.source_agent_id) == source_agent_id)
    if doc_type:
        stmt = stmt.where(col(Document.doc_type) == doc_type)
    if origin:
        stmt = stmt.where(col(Document.origin) == origin)
    if status_filter:
        stmt = stmt.where(col(Document.status) == status_filter)

    result = await session.exec(stmt)
    docs = result.all()
    return [await _resolve_agent(session, d, ctx.organization.id) for d in docs]


@router.post("", response_model=DocumentRead, status_code=status.HTTP_201_CREATED)
async def create_document(
    body: DocumentCreate,
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = ORG_MEMBER_DEP,
) -> DocumentRead:
    doc = Document(
        organization_id=ctx.organization.id,
        **body.model_dump(),
    )
    session.add(doc)
    await session.commit()
    await session.refresh(doc)
    return await _resolve_agent(session, doc, ctx.organization.id)


@router.get("/discover")
async def discover_documents_route(
    ctx: OrganizationContext = ORG_MEMBER_DEP,
) -> list[dict]:
    """Discover local files that can be imported."""
    if not settings.bridge_url:
        return []
    url = f"{settings.bridge_url.rstrip('/')}/documents/discover"
    headers = {"X-Bridge-Token": settings.bridge_token}
    async with httpx.AsyncClient(timeout=30.0) as client:
        resp = await client.get(url, headers=headers)
    if resp.status_code >= 400:
        return []
    return resp.json()


class ImportRequest(SQLModel):
    file_path: str
    title: str | None = None
    doc_type: str | None = None
    source_agent_id: UUID | None = None


@router.post("/import", response_model=DocumentRead, status_code=status.HTTP_201_CREATED)
async def import_document_route(
    body: ImportRequest,
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = ORG_MEMBER_DEP,
) -> DocumentRead:
    """Import a local file as a document record."""
    fpath = Path(body.file_path)
    mime = mimetypes.guess_type(str(fpath))[0] or "application/octet-stream"
    doc_type = body.doc_type
    if not doc_type:
        if "pdf" in mime: doc_type = "pdf"
        elif "presentation" in mime or "pptx" in mime: doc_type = "slide"
        elif "word" in mime or "docx" in mime: doc_type = "report"
        else: doc_type = "markdown"

    title = body.title or fpath.stem.replace("-", " ").replace("_", " ").title()
    doc = Document(
        organization_id=ctx.organization.id,
        title=title, doc_type=doc_type, source_agent_id=body.source_agent_id,
        file_path=body.file_path, mime_type=mime,
        file_size=fpath.stat().st_size if fpath.exists() else None,
        status="published",
    )
    session.add(doc)
    await session.commit()
    await session.refresh(doc)
    return await _resolve_agent(session, doc, ctx.organization.id)


class BatchImportResult(SQLModel):
    imported: int = 0
    skipped: int = 0
    errors: list[str] = []


@router.post("/batch-import", response_model=BatchImportResult)
async def batch_import_documents(
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = ORG_MEMBER_DEP,
) -> BatchImportResult:
    """Discover all local files and import any not already in the DB."""
    if not settings.bridge_url:
        raise HTTPException(status_code=503, detail="Bridge not configured")

    # Discover files
    url = f"{settings.bridge_url.rstrip('/')}/documents/discover"
    headers = {"X-Bridge-Token": settings.bridge_token}
    async with httpx.AsyncClient(timeout=60.0) as client:
        resp = await client.get(url, headers=headers)
    if resp.status_code >= 400:
        raise HTTPException(status_code=502, detail="Discovery failed")

    discovered = resp.json()

    # Get existing file paths to deduplicate
    existing_stmt = select(Document.file_path).where(
        col(Document.organization_id) == ctx.organization.id,
        col(Document.file_path).isnot(None),
    )
    result = await session.exec(existing_stmt)
    existing_paths = set(result.all())

    imported = 0
    skipped = 0
    errors: list[str] = []

    for f in discovered:
        fpath = f.get("path", "")
        if fpath in existing_paths:
            skipped += 1
            continue

        mime = f.get("mime_type", "application/octet-stream")
        doc_type = "markdown"
        if "pdf" in mime: doc_type = "pdf"
        elif "presentation" in mime or "pptx" in mime: doc_type = "slide"
        elif "word" in mime or "docx" in mime: doc_type = "report"
        elif "spreadsheet" in mime or "xlsx" in mime: doc_type = "report"

        try:
            doc = Document(
                organization_id=ctx.organization.id,
                title=f.get("name", Path(fpath).stem),
                doc_type=doc_type,
                file_path=fpath,
                mime_type=mime,
                file_size=f.get("size"),
                status="published",
            )
            session.add(doc)
            imported += 1
        except Exception as exc:
            errors.append(f"{fpath}: {str(exc)[:100]}")

    await session.commit()
    return BatchImportResult(imported=imported, skipped=skipped, errors=errors)


@router.get("/{doc_id}", response_model=DocumentRead)
async def get_document(
    doc_id: UUID,
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = ORG_MEMBER_DEP,
) -> DocumentRead:
    doc = await Document.objects.filter_by(
        id=doc_id,
        organization_id=ctx.organization.id,
    ).first(session)
    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND)
    return await _resolve_agent(session, doc, ctx.organization.id)


@router.patch("/{doc_id}", response_model=DocumentRead)
async def update_document(
    doc_id: UUID,
    body: DocumentUpdate,
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = ORG_MEMBER_DEP,
) -> DocumentRead:
    doc = await Document.objects.filter_by(
        id=doc_id,
        organization_id=ctx.organization.id,
    ).first(session)
    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND)

    for key, val in body.model_dump(exclude_unset=True).items():
        setattr(doc, key, val)
    doc.updated_at = utcnow()

    session.add(doc)
    await session.commit()
    await session.refresh(doc)
    return await _resolve_agent(session, doc, ctx.organization.id)


@router.get("/{doc_id}/preview")
async def preview_document(
    doc_id: UUID,
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = ORG_MEMBER_DEP,
) -> dict:
    """Extract readable content from a document for inline preview.

    Returns { html: str, text: str, preview_type: str } where:
    - html: rendered HTML for rich display (Word/PPT/Excel)
    - text: plain text fallback
    - preview_type: "html" | "text" | "unsupported"
    """
    doc = await Document.objects.filter_by(
        id=doc_id, organization_id=ctx.organization.id,
    ).first(session)
    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND)

    # If it has inline content (markdown), return that
    if doc.content and not doc.file_path:
        return {"html": "", "text": doc.content, "preview_type": "text"}

    # Need to fetch file from bridge and extract content
    if not doc.file_path or not settings.bridge_url:
        return {"html": "", "text": "", "preview_type": "unsupported"}

    # Fetch file bytes from bridge
    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            resp = await client.get(
                f"{settings.bridge_url.rstrip('/')}/documents/{doc_id}/download",
                headers={"X-Bridge-Token": settings.bridge_token},
            )
        if resp.status_code != 200:
            return {"html": "", "text": "Could not fetch file from bridge", "preview_type": "text"}
        file_bytes = resp.content
    except Exception as exc:
        return {"html": "", "text": f"Error fetching file: {str(exc)[:200]}", "preview_type": "text"}

    mime = doc.mime_type or ""
    title = (doc.title or "").lower()

    # Word documents (.docx)
    if "wordprocessingml" in mime or title.endswith(".docx"):
        try:
            import io
            from docx import Document as DocxDocument
            docx_doc = DocxDocument(io.BytesIO(file_bytes))
            html_parts = []
            for para in docx_doc.paragraphs:
                style = para.style.name if para.style else ""
                text = para.text.strip()
                if not text:
                    html_parts.append("<br/>")
                    continue
                if "Heading 1" in style:
                    html_parts.append(f"<h1>{text}</h1>")
                elif "Heading 2" in style:
                    html_parts.append(f"<h2>{text}</h2>")
                elif "Heading 3" in style:
                    html_parts.append(f"<h3>{text}</h3>")
                else:
                    # Check for bold/italic runs
                    run_html = ""
                    for run in para.runs:
                        t = run.text
                        if run.bold:
                            t = f"<strong>{t}</strong>"
                        if run.italic:
                            t = f"<em>{t}</em>"
                        run_html += t
                    html_parts.append(f"<p>{run_html or text}</p>")
            # Tables
            for table in docx_doc.tables:
                html_parts.append("<table border='1' cellpadding='4' cellspacing='0' style='border-collapse:collapse;width:100%'>")
                for row in table.rows:
                    html_parts.append("<tr>")
                    for cell in row.cells:
                        html_parts.append(f"<td>{cell.text}</td>")
                    html_parts.append("</tr>")
                html_parts.append("</table>")
            return {"html": "\n".join(html_parts), "text": "\n".join(p.text for p in docx_doc.paragraphs), "preview_type": "html"}
        except Exception as exc:
            return {"html": "", "text": f"Error parsing .docx: {str(exc)[:200]}", "preview_type": "text"}

    # Old .doc format — can't parse natively, return text extraction attempt
    if "msword" in mime or title.endswith(".doc"):
        try:
            text = file_bytes.decode("utf-8", errors="replace")
            # Try to extract readable text from binary .doc
            import re
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
            # Filter to lines with mostly printable chars
            lines = [l for l in text.split('\n') if len(l) > 2 and sum(c.isalpha() or c.isspace() for c in l) / max(len(l), 1) > 0.5]
            return {"html": "", "text": "\n".join(lines[:200]) if lines else "Unable to extract text from .doc format", "preview_type": "text"}
        except Exception:
            return {"html": "", "text": "Unable to extract text from .doc format", "preview_type": "text"}

    # Excel (.xlsx)
    if "spreadsheetml" in mime or title.endswith(".xlsx"):
        try:
            import io
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            html_parts = []
            for sheet_name in wb.sheetnames[:5]:  # Max 5 sheets
                ws = wb[sheet_name]
                html_parts.append(f"<h3>{sheet_name}</h3>")
                html_parts.append("<table border='1' cellpadding='4' cellspacing='0' style='border-collapse:collapse;width:100%'>")
                for row in ws.iter_rows(max_row=100, values_only=True):  # Max 100 rows
                    html_parts.append("<tr>")
                    for cell in row:
                        html_parts.append(f"<td>{cell if cell is not None else ''}</td>")
                    html_parts.append("</tr>")
                html_parts.append("</table>")
            wb.close()
            return {"html": "\n".join(html_parts), "text": "", "preview_type": "html"}
        except Exception as exc:
            return {"html": "", "text": f"Error parsing .xlsx: {str(exc)[:200]}", "preview_type": "text"}

    # PowerPoint (.pptx)
    if "presentationml" in mime or title.endswith(".pptx"):
        try:
            import io
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            html_parts = []
            for i, slide in enumerate(prs.slides):
                html_parts.append(f"<h3>Slide {i + 1}</h3>")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                html_parts.append(f"<p>{text}</p>")
            return {"html": "\n".join(html_parts), "text": "\n".join(p for p in html_parts if not p.startswith("<h")), "preview_type": "html"}
        except Exception as exc:
            return {"html": "", "text": f"Error parsing .pptx: {str(exc)[:200]}", "preview_type": "text"}

    # CSV / plain text
    if mime.startswith("text/") or "csv" in mime or "json" in mime or "xml" in mime:
        try:
            text = file_bytes.decode("utf-8", errors="replace")
            return {"html": "", "text": text[:50000], "preview_type": "text"}
        except Exception:
            pass

    return {"html": "", "text": "", "preview_type": "unsupported"}


@router.delete("/{doc_id}", status_code=status.HTTP_204_NO_CONTENT)
async def archive_document(
    doc_id: UUID,
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = Depends(require_org_admin),
) -> None:
    doc = await Document.objects.filter_by(
        id=doc_id,
        organization_id=ctx.organization.id,
    ).first(session)
    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND)
    doc.status = "archived"
    doc.updated_at = utcnow()
    session.add(doc)
    await session.commit()


@router.get("/{doc_id}/download")
async def download_document(
    doc_id: UUID,
    session: AsyncSession = SESSION_DEP,
    ctx: OrganizationContext = ORG_MEMBER_DEP,
):
    """Download/serve a file-backed document."""
    from fastapi.responses import StreamingResponse

    doc = await Document.objects.filter_by(
        id=doc_id,
        organization_id=ctx.organization.id,
    ).first(session)
    if not doc or not doc.file_path:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File-backed document not found")

    if not settings.bridge_url:
        raise HTTPException(status_code=503, detail="Bridge not configured")

    url = f"{settings.bridge_url.rstrip('/')}/documents/serve?path={doc.file_path}"
    headers = {"X-Bridge-Token": settings.bridge_token}
    async with httpx.AsyncClient(timeout=60.0) as client:
        resp = await client.get(url, headers=headers)
    if resp.status_code >= 400:
        raise HTTPException(status_code=resp.status_code, detail="File not found on bridge")

    return StreamingResponse(
        iter([resp.content]),
        media_type=doc.mime_type or "application/octet-stream",
        headers={"Content-Disposition": f'attachment; filename="{Path(doc.file_path).name}"'},
    )
